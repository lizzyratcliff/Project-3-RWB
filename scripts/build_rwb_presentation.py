"""
Build KIN 7518 group deck: outline structure + RWB content, RQ1 & RQ2, colors, speaker notes.

Usage (repo root):
  python scripts/build_rwb_presentation.py

Output: RWB_KIN7518_Presentation.pptx
Requires: python-pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
OUT = REPO / "RWB_KIN7518_Presentation.pptx"
VISUAL = REPO / "RWB_VISUAL.png"

# --- Theme (high-contrast, presentation-friendly) ---
NAVY = RGBColor(15, 30, 75)
INDIGO = RGBColor(67, 56, 202)
TEAL = RGBColor(13, 148, 136)
CORAL = RGBColor(225, 85, 84)
GOLD = RGBColor(217, 119, 6)
CREAM = RGBColor(255, 251, 240)
SKY = RGBColor(238, 246, 255)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(40, 44, 52)
TITLE_SLIDE_BG = RGBColor(30, 58, 138)
QNA_BG = RGBColor(49, 46, 129)


def _rgb(fill, color: RGBColor) -> None:
    fill.solid()
    fill.fore_color.rgb = color


def _slide_bg(slide, color: RGBColor) -> None:
    _rgb(slide.background.fill, color)


def _left_accent(slide, prs: Presentation, color: RGBColor, width_in: float = 0.14) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(width_in),
        prs.slide_height,
    )
    _rgb(bar.fill, color)
    bar.line.fill.background()


def _notes(slide, text: str) -> None:
    ns = slide.notes_slide
    ns.notes_text_frame.text = text.strip()


def _style_title(tf, size_pt: int, color: RGBColor, bold: bool = True) -> None:
    for p in tf.paragraphs:
        p.font.size = Pt(size_pt)
        p.font.bold = bold
        p.font.color.rgb = color


def _bullet_slide(
    prs: Presentation,
    title: str,
    lines: list[str],
    *,
    bg: RGBColor,
    accent: RGBColor,
    title_color: RGBColor,
    body_color: RGBColor,
    body_size: int = 20,
    notes: str,
) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _slide_bg(slide, bg)
    _left_accent(slide, prs, accent)

    slide.shapes.title.text = title
    _style_title(slide.shapes.title.text_frame, 32, title_color)

    body = slide.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(body_size)
        p.font.color.rgb = body_color
        p.space_after = Pt(8)

    _notes(slide, notes)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _slide_bg(slide, TITLE_SLIDE_BG)
    slide.shapes.title.text = (
        "Conflict, Morality & Polarization\nin Sport–Politics Social Comments (B50)"
    )
    sub = slide.placeholders[1]
    sub.text = (
        "Group RWB — KIN 7518 (Social Issues in Sport)\n"
        "Elizabeth Ratcliff · Jardyn Washington · Isabelle Besselman\n"
        "April 13, 2026"
    )
    _style_title(slide.shapes.title.text_frame, 36, WHITE)
    for p in sub.text_frame.paragraphs:
        p.font.size = Pt(19)
        p.font.color.rgb = RGBColor(226, 232, 240)
    _notes(
        slide,
        """
 [Timing: about 30–45 seconds. Optional: each member say name + one role, e.g. data / methods / writing.]

        OPENING LINE (pick one): “Today we’re sharing what happens in comment sections when sports stories bump into politics.” / “We analyzed tens of thousands of public comments to see how moral and political language shows up across three big platforms.”

        Welcome everyone. We are Group RWB. This project is for KIN 7518 — Social Issues in Sport. We studied online comments on posts where sports and politics overlap, especially content that also drew a lot of audience talk related to Donald Trump. We did not pick Trump to endorse or attack anyone; the class dataset was built around that topical overlap.

        PLAIN-LANGUAGE HOOK: Imagine a highlight reel or athlete story where half the comments are about the game and half are about politicians. That thread is doing two jobs at once — sports talk and political identity talk. We care how that shows up in language: insults framed as morality (“traitor,” “shameful”), and cues that sort people into camps (clear pro- or anti- cues, or general partisan name-calling).

        WHAT WE DID NOT DO: We are not scoring people as “good” or “bad” voters. We are not proving that a platform caused anyone to think a certain way. We describe patterns in text and counts in one class corpus (B50).

        ROAD MAP: (1) Why this matters and our two research questions. (2) What data we have and how big each platform is. (3) How we measured “moralized” and “stance” and what extra we could do on X. (4) Main results for all platforms + chart. (5) Extra results for X only. (6) Big-picture takeaways and limits. (7) Questions.

        You will hear TWO research questions. RQ1 is the main question across Instagram, X, and YouTube. RQ2 zooms in on X only because that spreadsheet included whether an account had a blue checkmark and follower counts — so we can ask, descriptively, whether “who is commenting” lines up with engagement and with our topic labels. Everything is correlational / descriptive.

        IF THE ROOM LOOKS LOST: Offer a one-sentence version: “We counted how often comments use moral-war words and political keywords, then compared Instagram, X, and YouTube, and peeked at X commenters’ verification and follower levels.”
        """,
    )

    # --- Slide 2: Context + BOTH RQs ---
    _bullet_slide(
        prs,
        "Why this matters & our research questions",
        [
            "Social issue: Sport–politics crossover posts pull audiences into moral outrage and partisan cues — not just game talk.",
            "Context: Corpus B50 = comments on hybrid sport/politics media with substantial Trump-related audience discussion.",
            "RQ 1 (primary): How do moralized discourse & polarization cues differ across Instagram, X, and YouTube — and how do those labels line up with likes (and on X, retweets/replies)? Descriptive only.",
            "RQ 2 (X only): Among X comments, are blue-verified status or follower level tied to (a) engagement and (b) a higher share of moralized or stance-coded comments? Exploratory — X is our smallest slice (n=1,008).",
        ],
        bg=CREAM,
        accent=TEAL,
        title_color=TEAL,
        body_color=CHARCOAL,
        body_size=19,
        notes="""
        [Timing: about 2–3 minutes. This slide sets up the “problem” and both questions — don’t rush the last bullet.]

        SOCIAL ISSUE (everyday words): When sports and politics meet online — a headline that mixes athletes, elections, or polarizing public figures — comment sections often turn into mini-debates about who is corrupt, who deserves respect, who is dangerous, and which “side” is to blame. That is the course theme in one picture: conflict, morality, and polarization. It matters for fans’ experience, for journalists covering hybrid stories, and for leagues or teams whose social channels host that talk.

        DEFINE “POLARIZATION” FOR A GENERAL AUDIENCE: We mean visible cues in text that line up with political camps — for example clear pro- or anti-Trump phrases from our lists, or general partisan insults — not a full psychological measure of someone’s identity.

        DATA IN ONE SENTENCE: B50 is a class corpus — the same topical scrape across Instagram, X, and YouTube — so we are comparing platforms while holding the topic area roughly constant. We are not comparing “all of Instagram” to “all of YouTube” in the world; we are comparing comments collected in this project around sport–politics crossover with heavy Trump-related audience talk.

        WALK THE BULLETS:
        • Bullet 1 — Tie to real life: people experience this when a golf clip or athlete post becomes a thread about politics.
        • Bullet 2 — “Corpus B50” is just our label for the merged files; emphasize same topic across sites.
        • Bullet 3 — RQ1 in one breath: “Do moral and political-cue rates differ by platform, and do those labels correlate with likes (and on X, retweets/replies)?” Stress DESCRIPTIVE — likes can mean outrage, jokes, or algorithms, not endorsement.
        • Bullet 4 — RQ2 in one breath: “On X only, do checkmark status or follower size relate to engagement or to our moralized/stance shares?” Stress EXPLORATORY and SMALL N (~1,008).

        IF ASKED “WHY THREE PLATFORMS?”: Different norms — short posts vs long threads, different moderation, different who shows up — so language and engagement can look different even when the topic is similar.

        IF ASKED “WHY TRUMP-RELATED?”: That is where this dataset was collected; it is a slice of online talk, not a verdict on any politician.

        TRANSITION LINE: “Next we’ll show how big the dataset is on each site — the bar chart preview matters for reading every later number.”
        """,
    )

    # --- Slide 3: Data & scope ---
    _bullet_slide(
        prs,
        "Data & scope",
        [
            "Dataset: B50 — merged public comments from Instagram, X (Twitter), and YouTube class exports.",
            "Scale: 58,464 comments — YouTube 45,623 · Instagram 11,833 · X 1,008.",
            "Same topical scrape (sport–politics crossover; heavy Trump-related talk) so cross-platform comparisons are meaningful.",
            "Takeaway: YouTube dominates raw volume; we emphasize within-platform percentages so one huge site does not drown out the others.",
        ],
        bg=SKY,
        accent=INDIGO,
        title_color=INDIGO,
        body_color=CHARCOAL,
        notes="""
        [Timing: about 1 minute. Goal: inoculate the audience against misreading huge N differences later.]

        ETHICS / PRIVACY (one calm sentence): These were public comments collected for class; we report aggregates in the paper and slides, not a pile-on of individual users. The public GitHub repo has no raw spreadsheets — only methods, dictionaries, and summary tables.

        NUMBERS TO INTERNALIZE: Total 58,464 comments. YouTube 45,623 (~78% of rows). Instagram 11,833. X 1,008. If you only remember one thing: YouTube is the “big bar,” X is the “small bar.”

        WHY WITHIN-PLATFORM PERCENTAGES MATTER: If you compare raw counts across platforms, YouTube would “win” every theme just because there are more rows. So for RQ1 we emphasize “out of all YouTube comments in B50, what percent…?” vs “out of all Instagram comments…?” That is fairer than comparing 2,548 to 284 raw moralized hits without context.

        MERGE IN PLAIN ENGLISH: We stacked three exports and labeled each row with a platform variable. Same conceptual topic area; not necessarily the same single video cross-posted everywhere — but the class scrape is designed for cross-platform comparison of this theme.

        IF ASKED “IS THIS REPRESENTATIVE?”: No — it is one batch tied to specific posts, algorithms, and who chose to comment. We describe B50, not “all Americans” or “the whole platform.”

        IF ASKED “DATES?”: If your instructor wants collection window details, point them to the course data documentation; the slide focuses on scale and scope.

        TRANSITION: “Now how did we turn text into categories? That’s the methods slide.”
        """,
    )

    # --- Slide 4: Methods ---
    _bullet_slide(
        prs,
        "How we measured things (RQ1 & RQ 2)",
        [
            "Moralized (RQ 1): Binary yes/no if text matches any of four keyword families — virtue/vice, harm/care, fairness/cheating, loyalty/betrayal (see keywords_v1.txt).",
            "Stance / polarization cues (RQ 1): Categories like pro-Trump, anti-Trump, partisan_other, neutral_unclear, mixed — from phrase and token lists. Many real opinions stay “neutral_unclear” if they do not hit a list term — so polarized shares are conservative lower bounds.",
            "Stats (RQ 1): Descriptive tables; chi-square for platform × stance with Cramer’s V (effect size); Spearman correlation for moralized vs log(1+likes) within each platform.",
            "RQ 2 extras (X only): Split by blue_verified; split followers into low / mid / high thirds (tertiles). Report Wilson 95% intervals for key percentages — a better “margin of error” style band for proportions than a plain +/-.",
            "Validation: Manual double-coding spot-check (example: ~58% exact stance match vs automation on a 40-comment sample) — dictionaries are imperfect; sarcasm and slang trip us up.",
        ],
        bg=CREAM,
        accent=GOLD,
        title_color=NAVY,
        body_color=CHARCOAL,
        body_size=18,
        notes="""
        [Timing: about 2 minutes. This slide is dense — slow down; it prevents misinterpretation later.]

        BULLET 1 — MORALIZED (tell a story): “If a comment calls someone a traitor, says an election was rigged, or uses language about murder or danger as moral judgment, our ‘moralized’ flag may turn on — if it matches our family lists.” Four families: virtue/vice, harm/care, fairness/cheating, loyalty/betrayal. One hit in any family = moralized “yes.” It is a rough tool, not a philosopher’s full theory of morality.

        BULLET 2 — STANCE: We assign one primary stance label from keyword lists: pro-Trump, anti-Trump, partisan_other (partisan but not clearly pro/anti Trump in our rules), neutral_unclear, or mixed. CRITICAL PLAIN-ENGLISH LINE: “Neutral_unclear is not the same as moderate politics.” It often means “our list didn’t catch the wording.” So pro-Trump % and anti-Trump % are lower-bound snapshots of dictionary-visible cues, not full public opinion.

        BULLET 3 — STATS (RQ1): Chi-square tests whether platform and stance are associated in the big cross-tab. p-values get tiny with huge N, so we emphasize Cramer’s V (~0.09) — a small-to-modest effect size. Spearman correlation: we pair each comment’s moralized yes/no with log(1+likes) because likes are super skewed (most comments have very few). Rho tells you direction/magnitude of a monotonic pattern, not a causal effect.

        BULLET 4 — RQ2: “Blue_verified” is whatever the scrape recorded for a checkmark-style field. Tertiles: sort X authors by followers, split into three equal-sized groups so each third has about the same number of comments — low, mid, high reach within this file. Wilson 95% intervals: for percentages, think “plausible range given sample noise” — wider when subgroups are small.

        BULLET 5 — VALIDATION: Humans double-coded a random sample (~40 comments in our archived spot-check). Stance agreement with automation was about 57.5% exact match; moralized agreement higher (~87.5%). That gap is expected when sarcasm, emojis, or context matter. Use that honesty to build trust.

        IF ASKED “WHY KEYWORDS, NOT BERT / ChatGPT?”: For class we prioritized transparency and reproducibility — every term is in keywords_v1.txt. ML could be next step; it trades interpretability for flexibility.

        IF ASKED “LOG(1+LIKES)?”:“Plus one” avoids log of zero; log pulls in extreme viral outliers so one monster comment does not dominate the raw scale.

        TRANSITION: “With rules on the table, here is what we found across platforms — you’ll see the bar chart on the next slide.”
        """,
    )

    # --- Slide 5: RQ1 findings + chart ---
    try:
        s5 = prs.slides.add_slide(prs.slide_layouts[6])
    except Exception:
        s5 = prs.slides.add_slide(prs.slide_layouts[5])
    _slide_bg(s5, SKY)
    _left_accent(s5, prs, TEAL)

    tx = s5.shapes.add_textbox(Inches(0.35), Inches(0.32), Inches(9.0), Inches(0.85))
    tx.text_frame.text = "Key findings — RQ 1 (all platforms)"
    _style_title(tx.text_frame, 30, INDIGO)

    body = s5.shapes.add_textbox(Inches(0.35), Inches(1.15), Inches(6.15), Inches(5.9))
    tf = body.text_frame
    tf.word_wrap = True
    bullets_r1 = [
        "Volume (chart): YouTube >> Instagram >> X — always pair raw counts with within-platform %.",
        "Moralized rate: Instagram ~2.4% · X ~5.1% · YouTube ~5.6% of comments (within each site).",
        "Stance (dictionary-defined): Most rows neutral_unclear. YouTube shows the largest pro-Trump share (~10.5%) vs Instagram (~2.8%) & X (~5.8%). Platform × stance association is real but modest (Cramer’s V ~ 0.09).",
        "Engagement hint: On X, comments tagged pro-Trump have much higher mean likes than other stance buckets — could be controversy, timing, or visibility; not proof of agreement.",
    ]
    for i, b in enumerate(bullets_r1):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(17)
        p.font.color.rgb = CHARCOAL
        p.space_after = Pt(10)

    if VISUAL.exists():
        s5.shapes.add_picture(str(VISUAL), Inches(6.85), Inches(1.05), width=Inches(6.0))
    else:
        nb = s5.shapes.add_textbox(Inches(6.85), Inches(2.0), Inches(5.8), Inches(1.2))
        nb.text_frame.text = "(Chart: RWB_VISUAL.png)"
        for p in nb.text_frame.paragraphs:
            p.font.color.rgb = CORAL

    _notes(
        s5,
        """
        [Timing: about 2 minutes. Gesture at the chart often — it anchors the story.]

        CHART WALKTHROUGH: Y-axis is number of comments; X-axis is platform. Read the labels on the bars out loud: Instagram 11,833; X 1,008; YouTube 45,623. One sentence: “YouTube is not just a little bigger — it’s most of the rows in this dataset.”

        BULLET 1 — VOLUME: Pair with the chart. Emphasize that any story about “more moral words on YouTube” could partly be “more comments total,” so we lean on within-platform percentages from the results tables.

        BULLET 2 — MORALIZED RATES: Read the three percents slowly: Instagram ~2.4%, X ~5.1%, YouTube ~5.6%. Plain English: “In this scrape, explicit moral-war keyword hits are uncommon on Instagram and somewhat more common on X and YouTube — still single-digit percentages.” Avoid saying “people are more moral on YouTube” — we measured dictionary matches, not character.

        BULLET 3 — STANCE / POLARIZATION CUES: Most comments are neutral_unclear under strict lists — say that out loud so the room doesn’t think “most people have no politics.” Then contrast pro-Trump shares: YouTube ~10.5%, X ~5.8%, Instagram ~2.8% in this corpus. Cramer’s V ~0.09: “statistically detectable but modest — platform nudges the pattern; it doesn’t lock everyone into one stance.”

        BULLET 4 — ENGAGEMENT ON X: Mean likes for pro-Trump-tagged comments (~132 in our tables) are much higher than other stance buckets on average. IMMEDIATELY QUALIFY: could be a few viral comments, could be argument threads, could be visibility — “we are not saying likes mean everyone agreed.”

        PAUSE FOR CONFUSION: If someone looks skeptical about low moralized %, clarify: short comments like “LOL” or emoji-only won’t trigger long moral lists; many insults might still miss our terms.

        TEAM TIP: If two people present, one narrates chart + moralized, other narrates stance + engagement.

        TRANSITION: “RQ1 was all three sites. Next we isolate X because only that file had verification and follower data — that’s RQ2.”
        """,
    )

    # --- Slide 6: RQ2 findings ---
    _bullet_slide(
        prs,
        "Key findings — RQ 2 (X / Twitter only)",
        [
            "Sample: 1,008 X comments — interpret as exploratory; subgroup cells get small fast.",
            "Blue check: Not verified n=175 → ~8.0% moralized (95% Wilson CI ~4.8–13.0%). Verified n=833 → ~4.4% (CI ~3.2–6.1%). In this scrape, unverified rows show a higher moralized % — could be composition, topic, or chance; we do not claim verification “causes” tone.",
            "Stance by verification (row %): e.g. pro-Trump ~4% (not verified) vs ~6.1% (verified); anti-Trump higher in the small not-verified slice (~4% vs ~1.4%).",
            "Followers (tertiles, equal n each): moralized ~7.1% (low followers) → ~4.2% (mid) → ~3.9% (high) — suggestive pattern only; same caution as above.",
            "Bottom line: X metadata reminds us “who comments” may differ by status and reach; do not treat X in B50 as representative of all users.",
        ],
        bg=RGBColor(255, 245, 245),
        accent=CORAL,
        title_color=RGBColor(153, 27, 27),
        body_color=CHARCOAL,
        body_size=17,
        notes="""
        [Timing: about 2 minutes. Speak slowly — this is the easiest slide to over-interpret.]

        OPENING LINE: “Now we zoom in on X only, because only that export had consistent fields for verification and followers.”

        SAMPLE SIZE MANTRA: 1,008 comments total. That is fine for class, small for sweeping claims. Whenever you split into subgroups, some buckets get tiny — be humble.

        BLUE CHECK — PLAIN ENGLISH: We used the scraped “verified / not verified” field. Verification rules and meaning changed over time on X; we are reporting what the dataset says, not endorsing checkmarks as “truth” or “quality.”

        MORALIZED BY VERIFICATION: Not verified n=175, ~8.0% moralized, Wilson interval roughly 4.8% to 13.0%. Verified n=833, ~4.4%, interval roughly 3.2% to 6.1%. SAY THIS CAREFULLY: “In this batch, unverified rows show a higher point estimate for moralized hits — but the intervals overlap, and the not-verified group is small, so don’t turn this into a headline about ‘verified people are nicer.’” Offer alternative explanations: different people attracted into threads, different moderation, viral outliers, sampling noise.

        STANCE BY VERIFICATION: Walk the audience through “row percentages” — within each verification group, what share of comments fall in each stance bucket? Pro-Trump share is higher among verified in our table (~6.1% vs ~4.0%); anti-Trump share is higher in the not-verified slice (~4% vs ~1.4%) but the not-verified group is small — stress uncertainty.

        FOLLOWER TERTILES: “We cut the X sample into thirds by follower count — low, mid, high — so each tertile has about the same number of comments.” Moralized ~7.1% → ~4.2% → ~3.9% from low to high followers in this file. Plain English: “suggestive step-down; still not causal.” Maybe bigger accounts face different norms; maybe content mix differs; we cannot tell from this table alone.

        IF ASKED “SHOULD WE TRUST VERIFIED MORE?”: Our project does not rank whose opinion counts more — we describe patterns. Also recall selection: who comments on sport–politics crossover posts is not random.

        IF ASKED “WHAT IS WILSON?”: “A statistician-approved way to put a range around a percent when sample size is modest — better than ‘plus or minus’ from a headline poll.”

        CLOSE: “RQ2 is exploratory glue — it reminds us comment sections are not anonymous noise; who shows up can correlate with language and engagement in this scrape.”
        """,
    )

    # --- Slide 7: Interpretation ---
    _bullet_slide(
        prs,
        "Interpretation & implications",
        [
            "RQ 1: Hybrid sport–politics content pulls moral and partisan language; the mix differs by platform — YouTube carries more dictionary-visible pro-Trump cues in this scrape than Instagram, with X between on several metrics.",
            "RQ 2: On X, verification and follower tier relate descriptively to moralized share and stance mix — treat as exploratory; smallest groups need humility.",
            "For sport orgs & communicators: Expect political spillover in comments; set moderation norms and messaging knowing threads may polarize.",
            "Limits: Keyword methods miss sarcasm and implicit opinion; engagement is not agreement; B50 is not representative of all voters or all platforms.",
        ],
        bg=CREAM,
        accent=INDIGO,
        title_color=NAVY,
        body_color=CHARCOAL,
        body_size=18,
        notes="""
        [Timing: about 2 minutes. This is your “so what” — land it clearly, then list limits without sounding like you’re disowning the project.]

        TWO-SENTENCE SUMMARY: “Sport–politics crossover content pulls moral and political language into comment threads; in B50, the mix differs by platform. On X, account status and follower tiers nudge descriptive patterns — but X is a small slice, so we stay exploratory.”

        RQ1 TAKEAWAY: YouTube dominates volume; within-platform, YouTube shows the largest dictionary-visible pro-Trump share in this scrape; Instagram is lowest on that measure; X is intermediate on several comparisons. Moralized keyword hits are uncommon everywhere but not zero — higher on X/YouTube than Instagram here. Effect sizes for platform × stance are modest — nuance, not a simple “Instagram is civil, YouTube is not” story.

        RQ2 TAKEAWAY: Verification and follower groups differ in moralized rates and stance mixes in descriptive tables — interpret with humility and overlapping uncertainty for small groups. This supports media literacy: “typical commenter” is not one person.

        FOR SPORT ORGANIZATIONS / COMMS: Political spillover is a planning issue, not a surprise — hybrid posts can activate identity politics in replies. Moderation policies, pinned context, and how staff account replies are framed can steer tone. We are not prescribing law — we are flagging risk and variability.

        FOR EDUCATORS / STUDENTS: Good discussion prompt — “Why might the same topic produce different comment chemistry on different apps?” Tie to algorithms, norms, and who uses each platform.

        LIMITS CHECKLIST (say any three): (1) Keyword dictionaries miss sarcasm and implicit opinion. (2) Stance labels are conservative — lots stay neutral_unclear. (3) Engagement metrics are not votes of agreement. (4) One class corpus — not representative of all voters or all posts. (5) We analyze language present, not factual truth of claims.

        IF ASKED “POLICY FIX?”: Stay humble — “We describe patterns; policy involves values, law, and platform governance beyond our data.”

        IF ASKED “WHAT WOULD YOU DO NEXT?”: Bigger human-coded sample; richer context per thread; qualitative examples; or complementary methods — still with ethics and privacy.

        STRONG CLOSING LINE BEFORE Q&A: “Our goal was transparent measurement and honest limits — so people can interpret sport–politics comment spaces with clearer eyes.”

        If someone asks “so who is right politically?” — the project is not designed to answer that. We describe language patterns in a defined sample.
        """,
    )

    # --- Slide 8: Q&A ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _slide_bg(slide, QNA_BG)
    slide.shapes.title.text = "Thank you"
    slide.placeholders[1].text = "Questions?\n\ngithub.com/lizzyratcliff/Project-3-RWB"
    _style_title(slide.shapes.title.text_frame, 44, WHITE)
    for p in slide.placeholders[1].text_frame.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(221, 214, 254)
    _notes(
        slide,
        """
        [Timing: flexible. Thank the class and instructor by name if appropriate.]

        OPENING: “Thank you — happy to take questions. The GitHub link on screen has our plan, dictionaries, and scripts for transparency.”

        LIKELY QUESTIONS — SHORT ANSWER PROMPTS:
        • “Why Trump-related talk?” → Class corpus built around sport–politics crossover where that discourse was prevalent; not a random sports week.
        • “Does this prove platforms cause polarization?” → No — descriptive patterns only; many confounds (algorithms, who comments, post selection).
        • “Why so many neutral_unclear?” → Keyword stance coding is intentionally conservative; real opinion often doesn’t hit list phrases.
        • “Can AI do this better?” → Maybe for recall, but transparency suffers; we prioritized auditable rules + human spot-checks.
        • “Most surprising finding?” → (Team chooses) e.g. modest effect sizes despite huge N; or X verification pattern with overlapping CIs; or YouTube pro-Trump dictionary share vs Instagram in same topical scrape.
        • “Biggest limitation?” → Keywords + one corpus + public comments ≠ full attitudes; engagement ≠ agreement.
        • “Ethical issues?” → Public data still deserves care — aggregates first, avoid piling on individuals; don’t overgeneralize to groups.
        • “What about TikTok / Facebook?” → Not in B50; different platforms would need their own scrape and ethics review.

        IF STUMPED: “That’s a great question — we didn’t test that in this project, but here’s how we’d think about designing a follow-up…”

        HANDOFF: If multiple presenters, agree who fields methods vs theory vs data questions.

        END: “Thanks again — we’ll stay after if anyone wants to chat.”
        """,
    )

    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
